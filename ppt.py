from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL_TYPE
import os

# Footer function
def add_footer(slide, text="CloudSim Plus Advanced Simulation"):
    left = Inches(0.5)
    top = Inches(6.8)
    width = Inches(8)
    height = Inches(0.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(100, 100, 100)  # Light gray color
    p.alignment = PP_ALIGN.CENTER

# Slide with bullet points
def add_bullet_slide(prs, title, bullet_points, subtitle=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)  # Larger title font
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Dark blue color

    if subtitle:
        slide.placeholders[1].text = subtitle

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear default content
    tf.margin_top = Inches(0.2)  # Reduce top margin
    tf.margin_bottom = Inches(0.2)  # Reduce bottom margin

    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(16)  # Slightly smaller font size
        p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1  # Use theme accent color

    add_footer(slide)

# Add gradient background to the title slide
def add_gradient_background(slide, start_color, end_color):
    fill = slide.background.fill
    fill.gradient()  # Enable gradient fill
    fill.gradient_stops[0].color.rgb = start_color  # Start color
    fill.gradient_stops[1].color.rgb = end_color    # End color

# Create the presentation
def create_presentation():
    prs = Presentation()

    # --- Title Slide ---
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Simulation and Automation in Cloud Computing"
    slide.placeholders[1].text = "Using CloudSim and CI/CD Pipeline\nManan Bhimjiyani | 22070122107 | April 2025"
    
    # Add gradient background
    add_gradient_background(slide, RGBColor(0, 112, 192), RGBColor(0, 51, 102))  # Blue gradient
    
    # Style title and subtitle
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)  # Larger title font
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White color
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(20)
    slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White color

    add_footer(slide)

    # --- Abstract Slide ---
    add_bullet_slide(prs, "Abstract", [
        "This project explores cloud resource management simulation using CloudSim Plus and automates the process with a CI/CD pipeline.",
        "Key features include VM scheduling, load balancing, fault tolerance, and energy consumption tracking.",
        "Automation is achieved using GitHub Actions with OpenJDK 21."
    ])

    # --- Problem Definition Slide ---
    add_bullet_slide(prs, "Problem Definition", [
        "Challenges in cloud data centers:",
        "   - Overloaded resources leading to bottlenecks.",
        "   - Energy inefficiency due to suboptimal utilization.",
        "   - System failures causing service disruptions.",
        "   - Manual workflows hindering reproducibility.",
        "   - Limited observability and environment standardization."
    ])

    # --- Objectives Slide ---
    add_bullet_slide(prs, "Objectives", [
        "Simulate advanced cloud resource management using CloudSim Plus.",
        "Optimize cloudlet scheduling using round-robin load balancing.",
        "Improve fault tolerance by simulating host failures and VM migration.",
        "Monitor and measure energy efficiency through integrated power models.",
        "Implement DevOps practices including CI/CD, automated testing, and observability."
    ])

    # --- Implementation Slide ---
    add_bullet_slide(prs, "Implementation", [
        "The simulation is implemented in Java using CloudSim Plus 8.0.0 and OpenJDK 21.",
        "Key components:",
        "   - Datacenter: 5 hosts, each with 8 PEs (3000 MIPS), 64 GB RAM, and a power model (250W max, 50W idle).",
        "   - VMs: 10 VMs, each with 4 PEs (2000 MIPS), 16 GB RAM.",
        "   - Cloudlets: 20 tasks with lengths 10k–29k MI.",
        "   - Load balancing: Round-robin assignment of cloudlets to VMs.",
        "   - Fault tolerance: Host 0 fails at 5 seconds, triggering VM migration.",
        "   - Energy tracking: Power consumption calculated every second."
    ])

    # --- Results Slide ---
    add_bullet_slide(prs, "Results and Analysis", [
        "Power consumption dropped to 600.00 W at 5 seconds due to Host 0 failure.",
        "All 20 cloudlets completed successfully despite the failure.",
        "Total energy consumption: 0.17 Wh.",
        "Total simulation time: 14.71 seconds."
    ])

    # --- Chart Slide ---
    slide_layout = prs.slide_layouts[5]  # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Power Usage Over Time"

    if os.path.exists("power_usage_chart.png"):
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(4)
        slide.shapes.add_picture("power_usage_chart.png", left, top, height=height)
    else:
        left = Inches(2)
        top = Inches(2)
        width = Inches(6)
        height = Inches(2)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = "Chart Image Not Found.\nRun the simulation to generate 'power_usage_chart.png'."
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)  # Red color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add warning icon
        left_icon = Inches(4)
        top_icon = Inches(3.5)
        icon = slide.shapes.add_shape(MSO_SHAPE.WARNING, left_icon, top_icon, Inches(1), Inches(1))
        icon.fill.solid()
        icon.fill.fore_color.rgb = RGBColor(255, 192, 0)  # Orange color

    add_footer(slide)

    # --- Conclusion Slide ---
    add_bullet_slide(prs, "Conclusion", [
        "Successfully simulated and automated cloud resource management using CloudSim Plus.",
        "Achievements include optimized load balancing, enhanced fault tolerance, and improved energy efficiency.",
        "Integrated CI/CD pipeline streamlined development and testing.",
        "Future enhancements may include advanced algorithms and multi-datacenter simulations."
    ])

    # --- Save ---
    prs.save("CloudSimPlus_Presentation.pptx")
    print("✅ Presentation saved as CloudSimPlus_Presentation.pptx.")

if __name__ == "__main__":
    create_presentation()